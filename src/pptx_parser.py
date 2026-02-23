"""
Conversion de fichiers PowerPoint (.pptx) en Markdown.

Deux modes de fonctionnement :

- Avec vision (défaut) :
    Tente de convertir le PPTX en PDF via LibreOffice, puis traite chaque page
    avec le modèle de vision (même pipeline que le PDF). C'est la méthode qui
    donne la meilleure qualité : schémas, tableaux et organigrammes sont vus
    dans leur contexte visuel complet.
    Si LibreOffice n'est pas installé, repli automatique sur l'extraction par shapes.

- Sans vision (--no-images) :
    Extraction rapide du texte de chaque shape, triée par ordre de lecture
    (haut→bas, gauche→droite). Aucun appel au modèle.

LibreOffice (soffice) doit être installé pour la méthode vision.
Téléchargement : https://www.libreoffice.org/download/
"""

import subprocess
import tempfile
from io import BytesIO
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER

from .image_analyzer import analyze_image

# ── Détection de LibreOffice ───────────────────────────────────────────────────

_LIBREOFFICE_CANDIDATES = [
    "soffice",
    "soffice.exe",
    r"C:\Program Files\LibreOffice\program\soffice.exe",
    r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
    "/usr/bin/soffice",
    "/usr/lib/libreoffice/program/soffice",
    "/Applications/LibreOffice.app/Contents/MacOS/soffice",
]

# ── Extraction par shapes (fallback) ──────────────────────────────────────────

_TITLE_PLACEHOLDER_TYPES = {
    PP_PLACEHOLDER.TITLE,
    PP_PLACEHOLDER.CENTER_TITLE,
}

# Tolérance verticale (EMU) pour regrouper des shapes sur la même "ligne"
ROW_TOLERANCE_EMU = 50_000


def _find_libreoffice() -> str | None:
    """
    Recherche l'exécutable LibreOffice (soffice) dans les emplacements courants.
    Retourne son chemin si trouvé, None sinon.
    """
    for candidate in _LIBREOFFICE_CANDIDATES:
        try:
            result = subprocess.run(
                [candidate, "--version"],
                capture_output=True,
                timeout=10,
            )
            if result.returncode == 0:
                return candidate
        except (FileNotFoundError, subprocess.TimeoutExpired, OSError):
            continue
    return None


def _convert_pptx_to_pdf(pptx_path: str, soffice: str, tmp_dir: str) -> Path | None:
    """
    Convertit un fichier PPTX en PDF via LibreOffice.
    Retourne le chemin du PDF généré, ou None en cas d'échec.
    """
    try:
        result = subprocess.run(
            [
                soffice,
                "--headless",
                "--convert-to", "pdf",
                "--outdir", tmp_dir,
                pptx_path,
            ],
            capture_output=True,
            timeout=180,
        )
        if result.returncode == 0:
            pdf_path = Path(tmp_dir) / (Path(pptx_path).stem + ".pdf")
            if pdf_path.exists():
                return pdf_path
    except (subprocess.TimeoutExpired, OSError):
        pass
    return None


def _is_title_placeholder(shape) -> bool:
    try:
        return (
            shape.is_placeholder
            and shape.placeholder_format.type in _TITLE_PLACEHOLDER_TYPES
        )
    except Exception:
        return False


def _shape_reading_order_key(shape) -> tuple:
    """Trie les shapes par ordre de lecture visuel (haut→bas, gauche→droite)."""
    top = shape.top if shape.top is not None else 0
    left = shape.left if shape.left is not None else 0
    return (round(top / ROW_TOLERANCE_EMU), left)


def _shape_text_to_markdown(shape, is_title: bool = False) -> str:
    """Convertit le contenu texte d'un shape en Markdown."""
    if not shape.has_text_frame:
        return ""
    parts: list[str] = []
    for para in shape.text_frame.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        level = para.level
        if is_title:
            parts.append(f"\n## {text}\n")
        elif level == 0:
            parts.append(f"\n{text}\n")
        else:
            indent = "  " * (level - 1)
            parts.append(f"\n{indent}- {text}\n")
    return "".join(parts)


def _parse_pptx_by_shapes(
    file_path: str,
    vision_model: str,
    analyze_images: bool,
    verbose: bool,
) -> str:
    """
    Extraction par shapes (fallback sans LibreOffice).

    Chaque diapositive est traitée shape par shape dans l'ordre visuel.
    Les images bitmap sont analysées individuellement via le modèle de vision.
    Les éléments graphiques vectoriels (flèches, connecteurs) ne sont pas capturés —
    c'est la limite de cette méthode par rapport au rendu LibreOffice.
    """
    prs = Presentation(file_path)
    slides_md: list[str] = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        if verbose:
            print(f"  Diapositive {slide_num}/{len(prs.slides)} ...", end=" ", flush=True)

        parts: list[str] = []

        # ── Titre ─────────────────────────────────────────────────────────
        title_shape = slide.shapes.title
        if title_shape and title_shape.has_text_frame:
            title_text = title_shape.text.strip()
            if title_text:
                parts.append(f"\n## {title_text}\n")
        else:
            parts.append(f"\n## Diapositive {slide_num}\n")

        # ── Autres shapes triées par ordre de lecture ──────────────────────
        other_shapes = sorted(
            [s for s in slide.shapes if s is not title_shape],
            key=_shape_reading_order_key,
        )

        for shape in other_shapes:
            if shape.has_text_frame:
                md = _shape_text_to_markdown(
                    shape, is_title=_is_title_placeholder(shape)
                )
                if md:
                    parts.append(md)

            if analyze_images and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    pil_img = Image.open(BytesIO(shape.image.blob))
                    if verbose:
                        print(
                            f"\n    Image ({pil_img.width}x{pil_img.height}) ...",
                            end=" ",
                            flush=True,
                        )
                    result = analyze_image(pil_img, vision_model)
                    if result["relevant"]:
                        parts.append(f"\n> **[Visuel]** {result['description']}\n")
                        if verbose:
                            print("pertinente [OK]")
                    else:
                        if verbose:
                            print(f"ignorée — {result['reason']}")
                except Exception as e:
                    if verbose:
                        print(f"\n    Avertissement image : {e}")

        if verbose:
            print("ok")

        if parts:
            slides_md.append("".join(parts))

    return "\n\n---\n\n".join(slides_md)


def parse_pptx(
    file_path: str,
    vision_model: str,
    analyze_images: bool = True,
    verbose: bool = True,
) -> str:
    """
    Parse un fichier PowerPoint et retourne son contenu en Markdown.

    Avec vision : conversion LibreOffice → PDF → rendu page par page (recommandé).
    Sans vision  : extraction rapide du texte des shapes (--no-images).
    """
    # Afficher le nombre de diapositives
    if verbose:
        prs_tmp = Presentation(file_path)
        print(f"  Nombre de diapositives : {len(prs_tmp.slides)}")
        del prs_tmp

    if not analyze_images:
        return _parse_pptx_by_shapes(file_path, vision_model, False, verbose)

    # ── Mode vision : tentative de conversion via LibreOffice ─────────────
    soffice = _find_libreoffice()

    if soffice:
        if verbose:
            print(f"  LibreOffice détecté — conversion en PDF pour rendu visuel...")
        with tempfile.TemporaryDirectory() as tmp_dir:
            pdf_path = _convert_pptx_to_pdf(file_path, soffice, tmp_dir)
            if pdf_path:
                from .pdf_parser import parse_pdf
                return parse_pdf(
                    str(pdf_path),
                    vision_model=vision_model,
                    analyze_images=True,
                    verbose=verbose,
                )
            elif verbose:
                print("  Conversion LibreOffice échouée — repli sur extraction par shapes.")
    elif verbose:
        print(
            "  LibreOffice non trouvé — extraction par shapes "
            "(les schémas vectoriels ne seront pas capturés).\n"
            "  Pour une meilleure qualité, installez LibreOffice : "
            "https://www.libreoffice.org/download/"
        )

    # ── Fallback : extraction par shapes ──────────────────────────────────
    return _parse_pptx_by_shapes(file_path, vision_model, analyze_images, verbose)
